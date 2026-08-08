import pandas as pd
import os
import logging
from typing import Optional, Union, Dict, Any
from concurrent.futures import ThreadPoolExecutor
import tempfile
import asyncio

logger = logging.getLogger(__name__)
executor = ThreadPoolExecutor(max_workers=4)


class DataLoader:
    def __init__(self):
        self.supported_formats = {
            ".csv", ".xlsx", ".xls", ".json", ".txt",
            ".sql", ".gsheet", ".pdf", ".docx", ".pptx"
        }

    # ============================================
    # 🚀 MAIN LOAD FUNCTION
    # ============================================
    async def load_data(
        self,
        file_path: Optional[str] = None,
        file_type: Optional[str] = None,
        sql_conn=None,
        sheet_name: Optional[Union[str, int]] = 0,
        query: Optional[str] = None,
        gsheet_url: Optional[str] = None,
        gsheet_credentials: Optional[str] = None,
    ) -> Dict[str, Any]:

        try:
            ext = (file_type or os.path.splitext(file_path or "")[1]).lower()

            if ext not in self.supported_formats:
                raise ValueError(f"Unsupported file format: {ext}")

            df = await self._run_in_thread(
                self._load_sync,
                file_path,
                ext,
                sql_conn,
                sheet_name,
                query,
                gsheet_url,
                gsheet_credentials
            )

            # 🔥 STEP 1: VALIDATION
            self._validate_dataframe(df)

            # 🔥 STEP 2: SCHEMA DETECTION
            schema = self._detect_schema(df)

            return {
                "dataframe": df,
                "schema": schema
            }

        except Exception as e:
            logger.exception(f"[UPLOAD ERROR] {str(e)}")
            raise

    async def _run_in_thread(self, func, *args):
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(executor, func, *args)

    # ============================================
    # 📥 FILE LOADERS
    # ============================================
    def _load_sync(
        self,
        file_path,
        ext,
        sql_conn,
        sheet_name,
        query,
        gsheet_url,
        gsheet_credentials
    ) -> pd.DataFrame:

        if ext == ".csv":
            return pd.read_csv(file_path, low_memory=False)

        elif ext in [".xlsx", ".xls"]:
            return pd.read_excel(file_path, sheet_name=sheet_name, engine="openpyxl")

        elif ext == ".json":
            return pd.read_json(file_path)

        elif ext == ".txt":
            return pd.read_csv(file_path, sep=None, engine="python")

        elif ext == ".sql":
            if not sql_conn or not query:
                raise ValueError("SQL requires both connection and query")
            return pd.read_sql(query, sql_conn)

        elif ext == ".gsheet":
            import gspread
            gc = gspread.service_account(filename=gsheet_credentials)
            sheet = gc.open_by_url(gsheet_url).sheet1
            return pd.DataFrame(sheet.get_all_records())

        elif ext == ".pdf":
            from tabula import read_pdf
            dfs = read_pdf(file_path, pages="all", multiple_tables=True)
            return pd.concat(dfs, ignore_index=True) if dfs else pd.DataFrame()

        elif ext == ".docx":
            import docx
            doc = docx.Document(file_path)
            text_data = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            return pd.DataFrame({"text": text_data})

        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            text_data = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_data.append(shape.text.strip())

            return pd.DataFrame({"text": text_data})

        else:
            raise ValueError(f"Unsupported file format: {ext}")

    # ============================================
    # ✅ VALIDATION
    # ============================================
    def _validate_dataframe(self, df: pd.DataFrame):

        if df is None:
            raise ValueError("Data loading failed (df is None)")

        if df.empty:
            raise ValueError("Loaded dataset is empty")

        if df.shape[1] < 2:
            raise ValueError("Dataset must have at least 2 columns")

    # ============================================
    # 🧠 SCHEMA DETECTION (FLAGSHIP 🔥)
    # ============================================
    def _detect_schema(self, df: pd.DataFrame) -> Dict[str, Any]:

        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        categorical_cols = df.select_dtypes(include="object").columns.tolist()
        datetime_cols = df.select_dtypes(include=["datetime64", "datetime64[ns]"]).columns.tolist()

        id_cols = [col for col in df.columns if "id" in col.lower()]

        # 🔥 METRIC DETECTION
        metric_keywords = ["revenue", "sales", "amount", "price", "score", "count"]
        metric_cols = [
            col for col in numeric_cols
            if any(k in col.lower() for k in metric_keywords)
        ]

        # 🔥 TIME DETECTION
        time_cols = datetime_cols + [
            col for col in df.columns
            if "date" in col.lower() or "time" in col.lower()
        ]

        # 🔥 BUSINESS CONTEXT
        context = "general"

        if any(k in " ".join(df.columns).lower() for k in ["revenue", "sales"]):
            context = "sales"

        elif any(k in " ".join(df.columns).lower() for k in ["customer", "user"]):
            context = "customer"

        elif any(k in " ".join(df.columns).lower() for k in ["product", "inventory"]):
            context = "product"

        return {
            "numeric_columns": numeric_cols,
            "categorical_columns": categorical_cols,
            "datetime_columns": datetime_cols,
            "id_columns": id_cols,
            "metric_columns": metric_cols,
            "time_columns": time_cols,
            "business_context": context
        }


# ============================================
# ✅ PIPELINE WRAPPER (UPDATED)
# ============================================

data_loader = DataLoader()


def load_data(file):
    """
    Decision Intelligence wrapper
    Returns BOTH dataframe + schema
    """

    tmp_path = None

    try:
        if hasattr(file, "seek"):
            file.seek(0)

        file_bytes = file.read()

        if not file_bytes:
            raise ValueError("Uploaded file is empty")

        suffix = ".csv"
        filename = getattr(file, "name", None)

        if filename and "." in filename:
            suffix = os.path.splitext(filename)[1].lower()

        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name

        try:
            result = asyncio.run(
                data_loader.load_data(file_path=tmp_path)
            )
        except RuntimeError:
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            result = loop.run_until_complete(
                data_loader.load_data(file_path=tmp_path)
            )

        df = result["dataframe"]
        schema = result["schema"]

        return df, schema  # 🔥 KEY UPGRADE

    except Exception as e:
        logger.exception(f"[PIPELINE LOAD ERROR] {str(e)}")
        raise

    finally:
        try:
            if tmp_path and os.path.exists(tmp_path):
                os.remove(tmp_path)
        except Exception:
            pass